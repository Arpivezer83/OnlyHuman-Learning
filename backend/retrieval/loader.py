import os
from pptx import Presentation

def load_pptx(file_path):
    """
    Betölti a PPTX fájlt abszolút elérési úttal.
    """
    base_dir = os.path.dirname(__file__)
    full_path = os.path.join(base_dir, file_path)
    full_path = os.path.abspath(full_path)

    if not os.path.exists(full_path):
        raise FileNotFoundError(f"Nem található a fájl: {full_path}")

    prs = Presentation(full_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def split_text(text, max_chunk_size=500):
    """
    Feldarabolja a szöveget kisebb darabokra.
    """
    words = text.split()
    chunks = []
    current_chunk = []

    for word in words:
        current_chunk.append(word)
        if len(current_chunk) >= max_chunk_size:
            chunks.append(' '.join(current_chunk))
            current_chunk = []

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks
